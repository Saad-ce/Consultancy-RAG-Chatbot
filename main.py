from pathlib import Path
from pptx import Presentation
from docx import Document
from unstructured.partition.auto import partition
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import requests
import json
import hashlib

# -------------------- Extraction --------------------
def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
        return "\n".join(text)
    except Exception:
        return None

def extract_text_from_docx(file_path):
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception:
        return None

def extract_text_with_unstructured(file_path):
    try:
        elements = partition(file_path=file_path)
        return "\n".join([str(el) for el in elements])
    except Exception:
        return None

def extract_text(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext == ".pdf":
        text = extract_text_with_unstructured(file_path)
    else:
        text = None
    return text or extract_text_with_unstructured(file_path) or ""

# -------------------- Chunking --------------------
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1200,
    chunk_overlap=150,
    separators=["\n\n", "\n• ", "\n- ", "\n", ". ", " ", ""]
)

base_path = Path("Diss_doc")  # update if needed
all_chunks = []

folders = sorted([p for p in base_path.iterdir() if p.is_dir()], key=lambda p: p.name)
for folder in folders:
    files = sorted([f for f in folder.rglob("*") if f.is_file()], key=lambda f: str(f))
    for file in files:
        if file.suffix.lower() in [".pptx", ".docx", ".pdf"]:
            raw_text = extract_text(file)
            if not raw_text:
                continue
            chunks = text_splitter.split_text(raw_text)

            doc_type = file.suffix.lower().lstrip(".")
            source_rel = f"{folder.name}/{file.name}"
            stat = file.stat()
            mtime = int(stat.st_mtime)
            doc_hash = hashlib.md5((str(file.resolve()) + str(mtime)).encode("utf-8")).hexdigest()

            for idx, chunk in enumerate(chunks):
                if not chunk.strip():
                    continue
                all_chunks.append({
                    "content": chunk,
                    "source": source_rel,
                    "doc_id": doc_hash,
                    "doc_type": doc_type,
                    "folder": folder.name,
                    "filename": file.name,
                    "chunk_idx": idx,
                    "chunk_char_len": len(chunk),
                    "file_mtime": mtime,
                })

# -------------------- Embeddings & FAISS --------------------
model = SentenceTransformer("all-MiniLM-L6-v2")
texts = [c["content"] for c in all_chunks]
metas = [c for c in all_chunks]

embeddings = model.encode(texts, show_progress_bar=True, normalize_embeddings=True)
embeddings = np.asarray(embeddings, dtype="float32")

dim = embeddings.shape[1]
index = faiss.IndexIDMap2(faiss.IndexFlatIP(dim))
ids = np.arange(len(embeddings), dtype=np.int64)
index.add_with_ids(embeddings, ids)

def search(query, top_k=5):
    query_vec = model.encode([query], normalize_embeddings=True)
    query_vec = np.asarray(query_vec, dtype="float32")
    scores, indices = index.search(query_vec, top_k)
    for i, idx in enumerate(indices[0]):
        if idx == -1:
            continue
        m = metas[idx]
        print(f"\n--- Match {i+1} ---")
        print(f"Cosine Similarity: {scores[0][i]:.4f}")
        print(f"Source: {m['source']} | chunk #{m['chunk_idx']} | type={m['doc_type']}")
        print(f"Content:\n{texts[idx][:500]}...")

# -------------------- Persistence --------------------
persist_dir = Path("vector_store")
persist_dir.mkdir(parents=True, exist_ok=True)
faiss.write_index(index, str(persist_dir / "docs.index"))
with open(persist_dir / "metas.json", "w", encoding="utf-8") as f:
    json.dump(metas, f, ensure_ascii=False)
np.save(persist_dir / "texts.npy", np.array(texts, dtype=object))

# -------------------- LLM + Prompt --------------------
def query_mistral(prompt, model="mistral:latest"):
    r = requests.post(
        "http://localhost:11434/api/generate",
        json={"model": model, "prompt": prompt, "stream": False}
    )
    return r.json()["response"]

def build_prompt(
    query,
    context_chunks_with_sources,
    history_snippets=None,
    mode="initial",
    max_items=8
):
    context_blocks = []
    for i, (chunk, src) in enumerate(context_chunks_with_sources, 1):
        context_blocks.append(f"<doc id='{i}' src='[{src}]'>\n{chunk}\n</doc>")
    context_text = "\n".join(context_blocks) if context_blocks else "<doc id='0'>N/A</doc>"

    history_text = ""
    if history_snippets:
        joined = "\n- ".join(history_snippets[:6])
        history_text = f"--- RECENT DIALOGUE CONTEXT (summaries) ---\n- {joined}\n"

    if mode == "initial":
        instructions = f"""You are a senior consultant leading the discovery phase of a client project.

You are given excerpts from past project documents. These are your ONLY source of information. Do NOT use external knowledge.
Each excerpt is enclosed in <doc> with a src filename in [brackets]. Stay strictly within what’s written.

{history_text}--- START OF CONTEXT ---
{context_text}
--- END OF CONTEXT ---

Client request:
\"\"\"{query}\"\"\"

Instructions (follow exactly):
- Use ONLY the provided <doc> content.
- If a required field is not explicitly present, write "N/A".
- Return at most {max_items} items.

Output format (strictly this structure, no extra text):
- [filename]: short factual statement
  - technologies: <comma-separated or N/A>
  - industries: <comma-separated or N/A>
  - timeframe: <as stated or N/A>

Begin your answer now with the bullet list only."""
    else:
        instructions = f"""You are continuing an ongoing consulting discussion.
Resolve references like "you said" using the RECENT DIALOGUE CONTEXT first. Use document excerpts only when needed for specific facts. No outside knowledge.

{history_text}--- OPTIONAL FACT CONTEXT ---
{context_text}
--- END OPTIONAL FACT CONTEXT ---

Client request:
\"\"\"{query}\"\"\"

Instructions (follow exactly):
- Prioritise the recent dialogue; consult <doc> only if necessary to support a claim.
- Keep it brief and actionable. Do not restate the entire plan.
- Cite [filename] ONLY when you directly use a fact from a <doc>.
- Max 6 bullets OR 120 words, whichever is shorter.
- If nothing in dialogue or docs supports the request, say "N/A".

Output format:
- bullet points (concise). Include [filename] inline only when citing a fact."""
    return instructions

# -------------------- RAG --------------------
def _get_source_label(meta):
    if isinstance(meta, dict):
        return meta.get("source", meta.get("filename", "unknown"))
    return str(meta)

def rag(query, top_k=10, min_score=0.5, mode="initial", history_snippets=None):
    qv = model.encode([query], normalize_embeddings=True).astype(np.float32)
    scores, idxs = index.search(qv, top_k)

    pairs, ranked = [], []
    for rank, idx in enumerate(idxs[0]):
        if idx == -1:
            continue
        score = float(scores[0][rank])
        if score < min_score:
            continue
        chunk = texts[idx]
        src_label = _get_source_label(metas[idx])
        pairs.append((chunk, src_label))
        ranked.append((score, src_label))

    print("\n🔍 Matches Above Threshold:")
    if not pairs:
        print("Not enough info")
        return None

    for i, (s, src) in enumerate(ranked, 1):
        print(f"{i:>2}. {src}  |  score={s:.4f}")

    prompt = build_prompt(query, pairs, history_snippets=history_snippets, mode=mode)
    resp = query_mistral(prompt)

    print("\n🧠 LLM Response:")
    print(resp)
    return resp, pairs, ranked

# -------------------- Chatbot --------------------
def _summarize_for_history(text, max_len=220):
    t = " ".join(text.split())
    return t[:max_len] + ("…" if len(t) > max_len else "")

def run_chatbot(top_k_initial=30, top_k_followup=12, min_score=0.45):
    print("RAG chatbot ready. Type /quit to exit.")
    mode = "initial"
    history_snippets = []

    while True:
        user_q = input("\nYou: ").strip()
        if not user_q:
            continue
        if user_q.lower() in {"/quit", "exit", "quit"}:
            print("Bye.")
            break

        tk = top_k_initial if mode == "initial" else top_k_followup
        result = rag(
            user_q,
            top_k=tk,
            min_score=min_score,
            mode=mode,
            history_snippets=history_snippets if mode == "followup" else None,
        )
        if not result:
            continue
        resp, _, _ = result

        history_snippets.append("Client: " + _summarize_for_history(user_q))
        history_snippets.append("Consultant: " + _summarize_for_history(resp))
        mode = "followup"

# Start interactive loop
run_chatbot()
